from flask import Flask, request, send_file, render_template, jsonify
import os
from datetime import datetime
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from PIL import Image, ImageDraw
import io

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024

@app.route("/")
def index():
    return render_template("index.html")

def convert_pptx_to_pdf(pptx_path):
    """Converte PPTX para PDF"""
    try:
        prs = Presentation(pptx_path)
        
        pdf_buffer = io.BytesIO()
        pdf = canvas.Canvas(pdf_buffer, pagesize=letter)
        
        page_width, page_height = letter
        
        for slide_idx, slide in enumerate(prs.slides):
            # Criar imagem do slide
            img = Image.new('RGB', (1280, 720), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # Processar shapes
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Text frame
                        text = shape.text[:100]  # Primeiros 100 caracteres
                        draw.text((50, 50 + slide_idx * 20), text, fill=(0, 0, 0))
                except:
                    pass
            
            # Converter imagem para bytes
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            
            # Adicionar ao PDF
            img_reader = ImageReader(img_bytes)
            pdf.drawImage(img_reader, 0, 0, width=page_width, height=page_height)
            pdf.showPage()
        
        pdf.save()
        pdf_buffer.seek(0)
        return pdf_buffer
        
    except Exception as e:
        raise Exception(f"Erro na conversão: {str(e)}")

@app.route("/upload", methods=["POST"])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "Nenhum arquivo foi enviado"}), 400
        
        file = request.files["file"]
        
        if file.filename == '':
            return jsonify({"error": "Nenhum arquivo selecionado"}), 400
        
        if not file.filename.lower().endswith('.pptx'):
            return jsonify({"error": "Por favor, envie um arquivo .pptx"}), 400
        
        # Gerar nome único
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_')
        safe_filename = timestamp + file.filename.replace('.pptx', '')
        filepath = os.path.join(UPLOAD_FOLDER, safe_filename + '.pptx')
        
        # Salvar arquivo
        file.save(filepath)
        
        try:
            # Converter para PDF
            pdf_buffer = convert_pptx_to_pdf(filepath)
            pdf_path = os.path.join(UPLOAD_FOLDER, safe_filename + '.pdf')
            
            with open(pdf_path, 'wb') as f:
                f.write(pdf_buffer.getvalue())
            
            # Remover PPTX
            if os.path.exists(filepath):
                os.remove(filepath)
            
            # Enviar PDF
            return send_file(
                pdf_path,
                as_attachment=True,
                download_name=safe_filename + '.pdf',
                mimetype='application/pdf'
            )
        
        except Exception as e:
            if os.path.exists(filepath):
                os.remove(filepath)
            return jsonify({"error": f"Erro ao converter: {str(e)}"}), 500
    
    except Exception as e:
        return jsonify({"error": f"Erro ao processar: {str(e)}"}), 500

if __name__ == "__main__":
    app.run(debug=True)
